from pathlib import Path
import re
from pptx import Presentation

root = Path(__file__).resolve().parents[1]
extracted = root / 'Sulthan_KMM_extracted.txt'
pptx_path = root / 'Sulthan_KMM_presentation_draft.pptx'
if not extracted.exists():
    print('Extracted text not found:', extracted)
    raise SystemExit(1)
if not pptx_path.exists():
    print('PPTX not found:', pptx_path)
    raise SystemExit(1)

text = extracted.read_text(encoding='utf-8')
# find DAFTAR PUSTAKA section
m = re.search(r"== DAFTAR PUSTAKA\n(.*?)(\n== |\Z)", text, flags=re.S)
refs = []
if m:
    block = m.group(1).strip()
    # split by double newlines or single newlines that look like new paragraph
    parts = [p.strip() for p in re.split(r"\n\n+", block) if p.strip()]
    refs = parts
else:
    # fallback: take last 20 lines
    refs = text.splitlines()[-30:]

prs = Presentation(str(pptx_path))
# add references slide
layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[5]
slide = prs.slides.add_slide(layout)
if slide.shapes.title:
    slide.shapes.title.text = 'Daftar Pustaka'
# find body placeholder
body = None
for shape in slide.placeholders:
    try:
        if shape.is_placeholder and shape.placeholder_format.type == 1:
            body = shape
            break
    except Exception:
        pass
if body is None:
    for shape in slide.shapes:
        if shape.has_text_frame:
            body = shape
            break
if body is None:
    print('No text frame available on slide')
    raise SystemExit(1)

tf = body.text_frame
tf.clear()
# add references (limit to 10 to avoid overcrowding)
for i, r in enumerate(refs[:10]):
    r_line = ' '.join(r.split())
    if i == 0:
        tf.text = r_line
    else:
        p = tf.add_paragraph()
        p.text = r_line
        p.level = 0

prs.save(str(pptx_path))
print('Updated PPTX with references at', pptx_path)
