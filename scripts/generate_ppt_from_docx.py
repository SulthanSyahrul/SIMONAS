import os
import zipfile
import xml.etree.ElementTree as ET
import re
from pathlib import Path
import sys
import subprocess

# workspace root is parent of this script's parent
root = Path(__file__).resolve().parents[1]
print('Workspace root:', root)
# find the docx file
candidates = list(root.glob('*.docx'))
docx_file = None
for p in candidates:
    name = p.name.lower()
    if 'anhutikno' in name or 'sulthan' in name:
        docx_file = p
        break
if not docx_file and candidates:
    docx_file = candidates[0]

if not docx_file:
    print('DOCX file not found in', root)
    sys.exit(1)

print('Using', docx_file)

# extract text and headings
ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
sections = []
current = {'heading': 'Pendahuluan', 'paras': []}
with zipfile.ZipFile(docx_file, 'r') as z:
    xml = z.read('word/document.xml')
    root_xml = ET.fromstring(xml)
    body = root_xml.find('w:body', ns)
    for p in body.findall('w:p', ns):
        texts = [t.text for t in p.findall('.//w:t', ns) if t.text]
        para_text = ''.join(texts).strip()
        if not para_text:
            continue
        pPr = p.find('w:pPr', ns)
        style = None
        if pPr is not None:
            pStyle = pPr.find('w:pStyle', ns)
            if pStyle is not None:
                style = pStyle.get('{%s}val' % ns['w'])
        if style and style.lower().startswith('heading'):
            sections.append(current)
            current = {'heading': para_text, 'paras': []}
        else:
            current['paras'].append(para_text)
    sections.append(current)

# write extracted text for review
extracted_path = root / 'Sulthan_KMM_extracted.txt'
with open(extracted_path, 'w', encoding='utf-8') as f:
    f.write(f'Source: {docx_file.name}\n\n')
    for sec in sections:
        f.write('== ' + sec['heading'] + '\n')
        for para in sec['paras']:
            f.write(para + '\n\n')

print('Wrote extracted text to', extracted_path)

# Prepare to create PPTX using python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    print('python-pptx not found, installing...')
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    from pptx import Presentation
    from pptx.util import Inches, Pt

prs = Presentation()
# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
# Use filename as title if available
title.text = docx_file.stem
subtitle.text = 'Ringkasan Laporan KMM - Sulthan Syahrul Bunayya Anhutikno'

# For each section create a slide with bullets (limit to first 6 bullets)
bullet_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[5]
for sec in sections:
    heading = sec['heading']
    paras = sec['paras']
    if not heading and not paras:
        continue
    slide = prs.slides.add_slide(bullet_layout)
    if slide.shapes.title:
        slide.shapes.title.text = heading
    # find a text frame to use
    body_shape = None
    # prefer placeholder body
    for shape in slide.placeholders:
        if shape.is_placeholder and shape.placeholder_format.type == 1:
            body_shape = shape
            break
    if body_shape is None:
        # fallback: first shape with text_frame
        for shape in slide.shapes:
            if shape.has_text_frame:
                body_shape = shape
                break
    if body_shape is None:
        continue
    tf = body_shape.text_frame
    tf.clear()
    # build bullets from paragraphs: take first sentence fragments
    bullets = []
    for para in paras:
        sents = re.split(r'(?<=[.!?])\s+', para)
        for s in sents:
            s = s.strip()
            if s:
                bullets.append(s)
            if len(bullets) >= 6:
                break
        if len(bullets) >= 6:
            break
    if not bullets:
        bullets = paras[:3]
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

out_pptx = root / 'Sulthan_KMM_presentation_draft.pptx'
prs.save(out_pptx)
print('Wrote PPTX to', out_pptx)

print('\nCREATED_FILES:')
print(extracted_path)
print(out_pptx)
