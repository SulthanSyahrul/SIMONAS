# -*- coding: utf-8 -*-
"""
Generate SIMONAS TA Seminar Hasil Presentation (max 15 slides)
Sulthan Syahrul Bunayya Anhutikno - V3423082
D3 Teknik Informatika - Sekolah Vokasi - Universitas Sebelas Maret
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───────────────────────────────────────────────────────
PRIMARY       = RGBColor(0x1A, 0x56, 0xDB)  # Deep Blue
PRIMARY_DARK  = RGBColor(0x0F, 0x3D, 0xA8)  # Darker Blue
ACCENT        = RGBColor(0x00, 0xB4, 0xD8)  # Cyan accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x1E, 0x1E, 0x2E)
DARK_TEXT      = RGBColor(0x2D, 0x3A, 0x4A)
LIGHT_BG      = RGBColor(0xF0, 0xF4, 0xF8)
SUBTITLE_GRAY = RGBColor(0x64, 0x74, 0x8B)
SUCCESS_GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE        = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT    = RGBColor(0xEF, 0x44, 0x44)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_gradient_header(slide, color1=PRIMARY, color2=PRIMARY_DARK, height=Inches(1.6)):
    """Add a colored header bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_bottom_bar(slide, color=ACCENT, height=Inches(0.08)):
    """Add a thin accent bar at the bottom."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_HEIGHT - height, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_slide_number(slide, number, total=15):
    """Add slide number in bottom right."""
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.5), Inches(1.2), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = SUBTITLE_GRAY


def add_title_on_header(slide, title_text, subtitle_text=None):
    """Add title text on the header bar."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)


def add_body_text(slide, text, left=Inches(0.8), top=Inches(2.0), width=Inches(11.5), height=Inches(4.8), font_size=Pt(18)):
    """Add body text block."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(6)
    return tf


def add_bullet_items(slide, items, left=Inches(0.8), top=Inches(2.0), width=Inches(11.5), height=Inches(5.0), font_size=Pt(16), icon="▸"):
    """Add bulleted items with custom icon."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon}  {item}"
        p.font.size = font_size
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)
        p.line_spacing = Pt(font_size.pt * 1.5)
    
    return tf


def add_info_card(slide, title, content, left, top, width=Inches(5.2), height=Inches(2.2), color=PRIMARY):
    """Add an info card with colored left border."""
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    
    # Left color bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(0.2), width - Inches(0.5), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    # Content
    txBox2 = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(0.65), width - Inches(0.5), height - Inches(0.85))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_TEXT
    p2.line_spacing = Pt(20)


def add_role_card(slide, role_name, features, left, top, color, icon_text):
    """Add a role feature card."""
    width = Inches(2.7)
    height = Inches(3.8)
    
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1)
    
    # Icon circle
    icon_size = Inches(0.65)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        left + (width - icon_size) / 2,
        top + Inches(0.25),
        icon_size, icon_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # Icon text
    txIcon = slide.shapes.add_textbox(
        left + (width - icon_size) / 2,
        top + Inches(0.3),
        icon_size, icon_size
    )
    tfI = txIcon.text_frame
    tfI.paragraphs[0].text = icon_text
    tfI.paragraphs[0].font.size = Pt(22)
    tfI.paragraphs[0].font.color.rgb = WHITE
    tfI.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Role name
    txName = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(1.05), width - Inches(0.3), Inches(0.4))
    tfN = txName.text_frame
    tfN.paragraphs[0].text = role_name
    tfN.paragraphs[0].font.size = Pt(15)
    tfN.paragraphs[0].font.bold = True
    tfN.paragraphs[0].font.color.rgb = color
    tfN.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Features
    txFeatures = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.5), width - Inches(0.4), height - Inches(1.7))
    tfF = txFeatures.text_frame
    tfF.word_wrap = True
    for i, feat in enumerate(features):
        if i == 0:
            p = tfF.paragraphs[0]
        else:
            p = tfF.add_paragraph()
        p.text = f"• {feat}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(4)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ═══════════════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Full blue background
    bg_shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = PRIMARY
    bg_shape.line.fill.background()
    
    # Decorative circle (top right)
    circle1 = slide1.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = PRIMARY_DARK
    circle1.line.fill.background()
    circle1.fill.fore_color.brightness = 0.1
    
    # Decorative circle (bottom left)
    circle2 = slide1.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(5), Inches(5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = ACCENT
    circle2.line.fill.background()
    
    # Title label "SEMINAR HASIL TUGAS AKHIR"
    txLabel = slide1.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(0.5))
    tf = txLabel.text_frame
    p = tf.paragraphs[0]
    p.text = "SEMINAR HASIL TUGAS AKHIR"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Main title
    txTitle = slide1.shapes.add_textbox(Inches(1), Inches(1.9), Inches(10), Inches(2.5))
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PENGEMBANGAN APLIKASI SIMONAS BERBASIS FLUTTER DAN SUPABASE UNTUK PENGAWASAN KELAS DI SMP NEGERI 1 JENAR"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.line_spacing = Pt(42)
    
    # Divider line
    line = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.6), Inches(2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    # Author info
    txAuthor = slide1.shapes.add_textbox(Inches(1), Inches(4.9), Inches(10), Inches(2.0))
    tf = txAuthor.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sulthan Syahrul Bunayya Anhutikno"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "V3423082"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    
    p3 = tf.add_paragraph()
    p3.text = ""
    p3.font.size = Pt(8)
    
    p4 = tf.add_paragraph()
    p4.text = "D3 Teknik Informatika • Sekolah Vokasi • Universitas Sebelas Maret"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    
    p5 = tf.add_paragraph()
    p5.text = "Pembimbing: Sahirul Alim Tri Bawono, S.Kom., M.Eng."
    p5.font.size = Pt(14)
    p5.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    
    add_slide_number(slide1, 1)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: Daftar Isi / Outline
    # ═══════════════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide2)
    add_title_on_header(slide2, "Outline Presentasi")
    add_bottom_bar(slide2)
    
    outline_items = [
        "Latar Belakang & Rumusan Masalah",
        "Tujuan & Manfaat",
        "Metode Pengembangan (Prototyping)",
        "Arsitektur Sistem & Teknologi",
        "Perancangan Produk (Use Case, Activity, ERD)",
        "Implementasi Aplikasi SIMONAS",
        "Fitur Multi-Role (Kepala Sekolah, Guru, BK, Siswa)",
        "Demo Aplikasi",
        "Pengujian (Black-Box Testing)",
        "Kesimpulan & Saran",
    ]
    
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(outline_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {i+1:02d}   "
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = PRIMARY
        
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(17)
        run2.font.color.rgb = DARK_TEXT
        
        p.space_after = Pt(8)
    
    add_slide_number(slide2, 2)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: Latar Belakang
    # ═══════════════════════════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide3)
    add_title_on_header(slide3, "Latar Belakang", "Kondisi Existing di SMP Negeri 1 Jenar")
    add_bottom_bar(slide3)
    
    # Problem cards
    add_info_card(slide3, 
        "❌  Masalah: Proses Manual",
        "• Absensi kelas dicatat pada lembar fisik\n• Jurnal mengajar ditulis di buku manual\n• Rekap data membutuhkan waktu lama\n• Rawan kesalahan & duplikasi data",
        Inches(0.6), Inches(1.9), Inches(5.5), Inches(2.5), RED_ACCENT)
    
    add_info_card(slide3, 
        "❌  Dampak Ketidakterpaduan",
        "• Kepala sekolah sulit melakukan monitoring real-time\n• BK kesulitan memperoleh data akademik cepat\n• Pekerjaan rekap berulang oleh guru dan wali kelas\n• Tidak ada integrasi data lintas fungsi",
        Inches(6.8), Inches(1.9), Inches(5.8), Inches(2.5), ORANGE)
    
    add_info_card(slide3, 
        "💡  Solusi: Aplikasi SIMONAS",
        "Mengembangkan aplikasi mobile berbasis Flutter & Supabase yang mengintegrasikan absensi, jurnal, jadwal, nilai, tugas, dan administrasi pembelajaran dalam satu platform multi-role.",
        Inches(0.6), Inches(4.8), Inches(12), Inches(2.0), SUCCESS_GREEN)
    
    add_slide_number(slide3, 3)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4: Rumusan Masalah & Tujuan
    # ═══════════════════════════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide4)
    add_title_on_header(slide4, "Rumusan Masalah & Tujuan")
    add_bottom_bar(slide4)
    
    add_info_card(slide4,
        "🔍  Rumusan Masalah",
        "Bagaimana mengembangkan aplikasi mobile berbasis Flutter dan Supabase yang dapat membantu pengawasan kelas dan pengelolaan akademik di SMP Negeri 1 Jenar secara terintegrasi?",
        Inches(0.6), Inches(1.9), Inches(12), Inches(1.6), PRIMARY)
    
    # Tujuan section
    tujuan_items = [
        "Menyediakan autentikasi & pemilihan role pengguna (Kepala Sekolah, Guru, BK, Siswa)",
        "Menyediakan dashboard sesuai hak akses masing-masing role",
        "Mengintegrasikan jurnal pembelajaran dan absensi siswa dalam satu alur",
        "Menyediakan modul jadwal, tugas, nilai, raport, & administrasi pembelajaran",
        "Menyediakan monitoring kelas, jurnal, & administrasi bagi kepala sekolah",
        "Menyusun rancangan pengujian fungsional (black-box testing)"
    ]
    
    txBox = slide4.shapes.add_textbox(Inches(0.6), Inches(3.7), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯  Tujuan Pengembangan"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    add_bullet_items(slide4, tujuan_items, Inches(0.8), Inches(4.2), Inches(11.5), Inches(3.0), Pt(14), "✓")
    add_slide_number(slide4, 4)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 5: Metode Pengembangan
    # ═══════════════════════════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide5)
    add_title_on_header(slide5, "Metode Pengembangan", "Pendekatan Prototyping")
    add_bottom_bar(slide5)
    
    stages = [
        ("01", "Pengumpulan\nKebutuhan", "Observasi, wawancara,\nanalisis proses bisnis\ndi SMP N 1 Jenar", PRIMARY),
        ("02", "Perancangan\nCepat", "Use case, activity\ndiagram, ERD, desain\nantarmuka", ACCENT),
        ("03", "Implementasi\nPrototipe", "Flutter + Dart,\nSupabase (Auth, DB,\nStorage), Riverpod", SUCCESS_GREEN),
        ("04", "Evaluasi &\nPenyempurnaan", "Black-box testing,\nwidget test, validasi\ndengan pengguna", ORANGE),
    ]
    
    start_x = Inches(0.5)
    card_width = Inches(2.8)
    gap = Inches(0.3)
    
    for i, (num, title, desc, color) in enumerate(stages):
        x = start_x + i * (card_width + gap)
        y = Inches(2.2)
        
        # Number circle
        circle = slide5.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.0), y, Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        txNum = slide5.shapes.add_textbox(x + Inches(1.0), y + Inches(0.1), Inches(0.8), Inches(0.7))
        tfN = txNum.text_frame
        tfN.paragraphs[0].text = num
        tfN.paragraphs[0].font.size = Pt(22)
        tfN.paragraphs[0].font.bold = True
        tfN.paragraphs[0].font.color.rgb = WHITE
        tfN.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow (except last)
        if i < 3:
            arrow_x = x + card_width
            arrow = slide5.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x - Inches(0.1), y + Inches(0.2), Inches(0.5), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
            arrow.line.fill.background()
        
        # Title
        txT = slide5.shapes.add_textbox(x, y + Inches(1.0), card_width, Inches(0.8))
        tfT = txT.text_frame
        tfT.word_wrap = True
        tfT.paragraphs[0].text = title
        tfT.paragraphs[0].font.size = Pt(16)
        tfT.paragraphs[0].font.bold = True
        tfT.paragraphs[0].font.color.rgb = color
        tfT.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Description card
        desc_card = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(1.9), card_width, Inches(1.8)
        )
        desc_card.fill.solid()
        desc_card.fill.fore_color.rgb = LIGHT_BG
        desc_card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        desc_card.line.width = Pt(1)
        
        txD = slide5.shapes.add_textbox(x + Inches(0.15), y + Inches(2.05), card_width - Inches(0.3), Inches(1.5))
        tfD = txD.text_frame
        tfD.word_wrap = True
        tfD.paragraphs[0].text = desc
        tfD.paragraphs[0].font.size = Pt(13)
        tfD.paragraphs[0].font.color.rgb = DARK_TEXT
        tfD.paragraphs[0].alignment = PP_ALIGN.CENTER
        tfD.paragraphs[0].line_spacing = Pt(20)
    
    # Note at bottom
    txNote = slide5.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11), Inches(0.5))
    tf = txNote.text_frame
    p = tf.paragraphs[0]
    p.text = "📌 Prototyping dipilih karena kebutuhan sekolah perlu divalidasi melalui observasi, diskusi, dan penyempurnaan bertahap"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = SUBTITLE_GRAY
    
    add_slide_number(slide5, 5)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 6: Arsitektur Sistem & Teknologi
    # ═══════════════════════════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide6)
    add_title_on_header(slide6, "Arsitektur Sistem & Teknologi", "Client-Server Architecture")
    add_bottom_bar(slide6)
    
    # Client box
    client_box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5)
    )
    client_box.fill.solid()
    client_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    client_box.line.color.rgb = PRIMARY
    client_box.line.width = Pt(2)
    
    txClient = slide6.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(3.4), Inches(0.4))
    tf = txClient.text_frame
    tf.paragraphs[0].text = "📱 CLIENT (Flutter)"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY

    client_features = [
        "Flutter 3.35.7 + Dart 3.9.2",
        "Riverpod State Management",
        "Feature-based Architecture",
        "Multi-platform (Android/iOS)",
        "Folder: features/, services/,\nmodels/, providers/"
    ]
    add_bullet_items(slide6, client_features, Inches(0.8), Inches(2.6), Inches(3.2), Inches(3.5), Pt(13), "•")
    
    # Arrow
    arrow = slide6.shapes.add_shape(
        MSO_SHAPE.LEFT_RIGHT_ARROW, Inches(4.5), Inches(3.8), Inches(1.2), Inches(0.6)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()
    
    txArrow = slide6.shapes.add_textbox(Inches(4.4), Inches(4.5), Inches(1.4), Inches(0.3))
    tf = txArrow.text_frame
    tf.paragraphs[0].text = "REST API"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Server box
    server_box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(2.0), Inches(3.8), Inches(4.5)
    )
    server_box.fill.solid()
    server_box.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    server_box.line.color.rgb = SUCCESS_GREEN
    server_box.line.width = Pt(2)
    
    txServer = slide6.shapes.add_textbox(Inches(6.1), Inches(2.1), Inches(3.4), Inches(0.4))
    tf = txServer.text_frame
    tf.paragraphs[0].text = "☁️ BACKEND (Supabase)"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = SUCCESS_GREEN

    server_features = [
        "Supabase Auth (autentikasi)",
        "PostgreSQL Database",
        "Supabase Storage (file)",
        "Row-Level Security (RLS)",
        "Real-time subscriptions"
    ]
    add_bullet_items(slide6, server_features, Inches(6.2), Inches(2.6), Inches(3.3), Inches(3.5), Pt(13), "•")
    
    # Security box
    sec_box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(2.0), Inches(3.0), Inches(4.5)
    )
    sec_box.fill.solid()
    sec_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    sec_box.line.color.rgb = ORANGE
    sec_box.line.width = Pt(2)
    
    txSec = slide6.shapes.add_textbox(Inches(10.2), Inches(2.1), Inches(2.6), Inches(0.4))
    tf = txSec.text_frame
    tf.paragraphs[0].text = "🔒 KEAMANAN"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE

    sec_features = [
        "RLS per tabel",
        "Policy auth.uid()",
        "Role-based access",
        "Kepala sekolah:\nread-only",
        "Guru: own data\nonly"
    ]
    add_bullet_items(slide6, sec_features, Inches(10.3), Inches(2.6), Inches(2.5), Inches(3.5), Pt(12), "•")
    
    add_slide_number(slide6, 6)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 7: Perancangan - Use Case
    # ═══════════════════════════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide7)
    add_title_on_header(slide7, "Perancangan Produk", "Use Case Diagram & Penggolongan Pengguna")
    add_bottom_bar(slide7)
    
    roles_data = [
        ("Kepala Sekolah", [
            "Dashboard monitoring",
            "Monitoring kelas real-time",
            "Monitoring jurnal guru",
            "Verifikasi administrasi",
            "Manajemen tahun ajaran",
            "Pengelolaan guru & mapel"
        ], PRIMARY, "KS"),
        ("Guru", [
            "Jadwal mengajar",
            "Input jurnal & absensi",
            "Pengelolaan nilai",
            "Pengelolaan tugas",
            "Upload administrasi",
            "Raport siswa"
        ], SUCCESS_GREEN, "G"),
        ("BK", [
            "Manajemen siswa",
            "Pengaturan kelas",
            "Kenaikan kelas",
            "Import data siswa",
            "Monitoring presensi"
        ], ORANGE, "BK"),
        ("Siswa", [
            "Lihat jadwal pelajaran",
            "Lihat nilai & tugas",
            "Histori presensi",
            "Informasi akademik"
        ], ACCENT, "S"),
    ]
    
    for i, (role, features, color, icon) in enumerate(roles_data):
        add_role_card(slide7, role, features, 
                     Inches(0.4) + i * Inches(3.15), Inches(2.2), color, icon)
    
    add_slide_number(slide7, 7)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 8: Perancangan - ERD & Database
    # ═══════════════════════════════════════════════════════════════════
    slide8 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide8)
    add_title_on_header(slide8, "Desain Basis Data", "Entity Relationship Diagram (ERD)")
    add_bottom_bar(slide8)
    
    db_modules = [
        ("Modul Pengguna", "• Tabel users, user_roles\n• Multi-role per akun\n• Foreign key ke profiles", PRIMARY),
        ("Modul Jadwal & Jurnal", "• Tabel jadwal, jurnal_kbm\n• Tabel absensi terintegrasi\n• Relasi guru-kelas-mapel", ACCENT),
        ("Modul Penilaian", "• Tabel nilai (UTS, UAS)\n• Tabel tugas & pengumpulan\n• Relasi siswa-kelas-mapel", SUCCESS_GREEN),
        ("Modul Admin Pembelajaran", "• Silabus, Prota, Promes, RPP\n• Storage file di Supabase\n• Verifikasi kepala sekolah", ORANGE),
    ]
    
    for i, (title, content, color) in enumerate(db_modules):
        col = i % 2
        row = i // 2
        add_info_card(slide8, f"📊  {title}", content,
                     Inches(0.5) + col * Inches(6.2),
                     Inches(2.0) + row * Inches(2.5),
                     Inches(5.8), Inches(2.2), color)
    
    add_slide_number(slide8, 8)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 9: Implementasi - Login & Role
    # ═══════════════════════════════════════════════════════════════════
    slide9 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide9)
    add_title_on_header(slide9, "Implementasi: Login & Pemilihan Role", "Autentikasi Multi-Role dengan Supabase Auth")
    add_bottom_bar(slide9)
    
    add_info_card(slide9,
        "🔐  Halaman Login",
        "• Validasi input username/email & password\n• Autentikasi via Supabase Auth\n• Error handling & feedback visual\n• Perlindungan kredensial",
        Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5), PRIMARY)
    
    add_info_card(slide9,
        "👤  Pemilihan Role",
        "• Akun multi-role: satu akun bisa punya >1 role\n• Setelah login, muncul daftar role tersedia\n• Routing otomatis ke dashboard sesuai role\n• State dikelola oleh Riverpod provider",
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5), ACCENT)
    
    # Code snippet
    code_box = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(12), Inches(2.2)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    code_box.line.fill.background()
    
    code_text = """Widget? buildDashboardForRole({required String role, required String userId}) {
  switch (normalizeRole(role)) {
    case 'kepala_sekolah': return KepalaSekolahDashboardScreen();
    case 'guru':           return GuruDashboardScreen(guruUid: userId);
    case 'BK':             return DashboardKemahasiswaanScreen(kemahasiswaanUid: userId);
    case 'siswa':          return SiswaDashboardScreen(siswaUid: userId);
    default: return null;
  }
}"""
    txCode = slide9.shapes.add_textbox(Inches(0.8), Inches(4.95), Inches(11.5), Inches(2.0))
    tf = txCode.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xA5, 0xF3, 0xFC)
    p.font.name = "Consolas"
    
    add_slide_number(slide9, 9)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 10: Implementasi - Modul Guru
    # ═══════════════════════════════════════════════════════════════════
    slide10 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide10)
    add_title_on_header(slide10, "Implementasi: Modul Guru", "Dashboard, Jurnal & Absensi Terintegrasi")
    add_bottom_bar(slide10)
    
    add_info_card(slide10,
        "📋  Dashboard Guru",
        "• Menu shortcut: Jadwal, Jurnal, Nilai, Berkas\n• Dropdown tahun ajaran di AppBar\n• Informasi ringkas kegiatan hari ini\n• Navigasi cepat ke setiap modul",
        Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3), SUCCESS_GREEN)
    
    add_info_card(slide10,
        "✏️  Jurnal & Absensi Terintegrasi",
        "• Form jurnal: tanggal, kelas, mapel, materi, metode, catatan\n• Absensi terintegrasi (bukan halaman terpisah)\n• Daftar siswa auto-update sesuai kelas dipilih\n• Status per siswa: Hadir, Izin, Sakit, Alpa\n• Ringkasan statistik kehadiran real-time",
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.3), PRIMARY)
    
    add_info_card(slide10,
        "📊  Modul Nilai & Tugas",
        "• Input nilai UTS, UAS per siswa per mapel\n• Pengelolaan tugas: buat, edit, deadline\n• Siswa submit tugas via aplikasi\n• Raport per semester",
        Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.3), ACCENT)
    
    add_info_card(slide10,
        "📁  Administrasi Pembelajaran",
        "• Upload Silabus, Prota, Promes, RPP\n• File dikompresi sebelum upload ke Supabase Storage\n• Status verifikasi oleh kepala sekolah\n• Organisasi per tahun ajaran & semester",
        Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.3), ORANGE)
    
    add_slide_number(slide10, 10)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 11: Implementasi - Modul Kepala Sekolah
    # ═══════════════════════════════════════════════════════════════════
    slide11 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide11)
    add_title_on_header(slide11, "Implementasi: Modul Kepala Sekolah", "Monitoring & Manajemen Akademik")
    add_bottom_bar(slide11)
    
    add_info_card(slide11,
        "📊  Dashboard Monitoring",
        "• Ringkasan kehadiran guru harian\n• Status pengumpulan RPP per guru\n• Statistik KBM keseluruhan\n• Overview administrasi pembelajaran",
        Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3), PRIMARY)
    
    add_info_card(slide11,
        "👁️  Monitoring Kelas Real-Time",
        "• Melihat status KBM setiap kelas saat ini\n• Identifikasi kelas yang belum terisi jurnal\n• Data kehadiran guru per jadwal\n• Alert untuk kelas tanpa aktivitas",
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.3), RED_ACCENT)
    
    add_info_card(slide11,
        "📝  Monitoring Jurnal",
        "• Review jurnal mengajar seluruh guru\n• Filter per guru, kelas, tanggal\n• Verifikasi kelengkapan jurnal\n• Histori jurnal per periode",
        Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.3), ACCENT)
    
    add_info_card(slide11,
        "⚙️  Manajemen Akademik",
        "• Manajemen tahun ajaran aktif\n• Pengelolaan data guru & mata pelajaran\n• Pengaturan jadwal\n• Verifikasi dokumen administrasi",
        Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.3), SUCCESS_GREEN)
    
    add_slide_number(slide11, 11)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 12: Implementasi - Modul BK & Siswa
    # ═══════════════════════════════════════════════════════════════════
    slide12 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide12)
    add_title_on_header(slide12, "Implementasi: Modul BK & Siswa")
    add_bottom_bar(slide12)
    
    add_info_card(slide12,
        "🏫  Modul BK (Bimbingan Konseling)",
        "• Manajemen data siswa (CRUD)\n• Import data siswa dari file CSV\n• Pengaturan kelas & penempatan siswa\n• Proses kenaikan kelas\n• Monitoring presensi siswa lintas kelas",
        Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.2), ORANGE)
    
    add_info_card(slide12,
        "🎒  Modul Siswa",
        "• Dashboard informasi akademik personal\n• Lihat jadwal pelajaran harian\n• Akses nilai UTS/UAS per mata pelajaran\n• Melihat & mengumpulkan tugas\n• Histori presensi (hadir, izin, sakit, alpa)\n• Akumulasi kehadiran per semester",
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(3.2), ACCENT)
    
    # Key highlight
    highlight = slide12.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(12), Inches(1.3)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    highlight.line.color.rgb = PRIMARY
    highlight.line.width = Pt(2)
    
    txHL = slide12.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.9))
    tf = txHL.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 Keunggulan: Semua role terintegrasi dalam satu aplikasi dengan single codebase (Flutter). Data tersinkronisasi secara real-time melalui Supabase."
    p.font.size = Pt(15)
    p.font.color.rgb = PRIMARY
    p.font.bold = True
    
    add_slide_number(slide12, 12)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 13: Demo Aplikasi (placeholder for live demo)
    # ═══════════════════════════════════════════════════════════════════
    slide13 = prs.slides.add_slide(blank_layout)
    
    # Full gradient bg
    bg = slide13.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    bg.line.fill.background()
    
    # Decorative circles
    c1 = slide13.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(7), Inches(7))
    c1.fill.solid()
    c1.fill.fore_color.rgb = PRIMARY_DARK
    c1.line.fill.background()
    
    c2 = slide13.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0x06, 0x4E, 0x3B)
    c2.line.fill.background()
    
    # Demo text
    txDemo = slide13.shapes.add_textbox(Inches(2), Inches(2.0), Inches(9), Inches(1.5))
    tf = txDemo.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🖥️  DEMO APLIKASI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Demo items
    demo_items = [
        "Login & Pemilihan Role",
        "Dashboard Guru → Jurnal & Absensi",
        "Dashboard Kepala Sekolah → Monitoring Kelas",
        "Dashboard BK → Manajemen Siswa",
        "Dashboard Siswa → Informasi Akademik"
    ]
    
    txItems = slide13.shapes.add_textbox(Inches(3), Inches(3.8), Inches(7), Inches(3.0))
    tf = txItems.text_frame
    tf.word_wrap = True
    for i, item in enumerate(demo_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
    
    add_slide_number(slide13, 13)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 14: Pengujian
    # ═══════════════════════════════════════════════════════════════════
    slide14 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide14)
    add_title_on_header(slide14, "Pengujian", "Black-Box Testing & Automated Testing")
    add_bottom_bar(slide14)
    
    add_info_card(slide14,
        "🧪  Metode Black-Box Testing",
        "• Fokus pada pemeriksaan input-output tanpa melihat kode internal\n• Skenario berdasarkan use case pengguna\n• Equivalence partitioning untuk setiap fitur\n• Validasi kesesuaian fungsi terhadap kebutuhan",
        Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5), PRIMARY)
    
    add_info_card(slide14,
        "🤖  Automated Testing",
        "• Widget test via flutter test\n• Verifikasi komponen login: judul, kolom input, tombol\n• Build test: flutter build bundle → 100% sukses\n• Zero syntax error, zero dependency failure",
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5), SUCCESS_GREEN)
    
    # Results summary
    result_box = slide14.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(12), Inches(2.2)
    )
    result_box.fill.solid()
    result_box.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    result_box.line.color.rgb = SUCCESS_GREEN
    result_box.line.width = Pt(2)
    
    txResult = slide14.shapes.add_textbox(Inches(0.8), Inches(4.95), Inches(11.5), Inches(2.0))
    tf = txResult.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✅  HASIL PENGUJIAN"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    
    p2 = tf.add_paragraph()
    p2.text = ""
    p2.font.size = Pt(6)
    
    results = [
        "Seluruh skenario black-box testing → BERHASIL (100%)",
        "Widget test: All tests passed!",
        "Build bundle: Kompilasi 100% sukses tanpa error",
        "Fungsionalitas utama berjalan sesuai spesifikasi kebutuhan"
    ]
    for r in results:
        p3 = tf.add_paragraph()
        p3.text = f"   ✓  {r}"
        p3.font.size = Pt(15)
        p3.font.color.rgb = DARK_TEXT
        p3.space_after = Pt(6)
    
    add_slide_number(slide14, 14)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 15: Kesimpulan & Saran
    # ═══════════════════════════════════════════════════════════════════
    slide15 = prs.slides.add_slide(blank_layout)
    add_gradient_header(slide15)
    add_title_on_header(slide15, "Kesimpulan & Saran")
    add_bottom_bar(slide15)
    
    # Kesimpulan
    conclusions = [
        "Aplikasi SIMONAS berhasil dikembangkan menggunakan Flutter + Supabase untuk digitalisasi akademik",
        "Sistem multi-role berfungsi dengan baik (Kepala Sekolah, Guru, BK, Siswa)",
        "Integrasi jurnal & absensi dalam satu alur mengurangi kerja berulang guru",
        "Monitoring real-time membantu kepala sekolah mengawasi KBM lebih efektif",
        "Pengujian black-box 100% berhasil → siap diimplementasikan"
    ]
    
    txK = slide15.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.2), Inches(0.4))
    tf = txK.text_frame
    tf.paragraphs[0].text = "📌  Kesimpulan"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY
    
    add_bullet_items(slide15, conclusions, Inches(0.6), Inches(2.3), Inches(6.2), Inches(4.5), Pt(13), "✓")
    
    # Saran
    suggestions = [
        "Pilot testing pada 1-2 kelas sebelum implementasi menyeluruh",
        "Penambahan fitur push notification untuk pengingat otomatis",
        "Rekapitulasi data dalam format PDF/Excel yang dapat diunduh",
        "Pengembangan modul orang tua/wali untuk monitoring anak",
        "Integrasi dengan sistem eksternal (e-Rapor)"
    ]
    
    txS = slide15.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.4))
    tf = txS.text_frame
    tf.paragraphs[0].text = "💡  Saran Pengembangan"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    
    add_bullet_items(slide15, suggestions, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.5), Pt(13), "▸")
    
    # Thank you bar
    thank_box = slide15.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(6.2), Inches(8), Inches(0.9)
    )
    thank_box.fill.solid()
    thank_box.fill.fore_color.rgb = PRIMARY
    thank_box.line.fill.background()
    
    txThank = slide15.shapes.add_textbox(Inches(2.5), Inches(6.3), Inches(8), Inches(0.7))
    tf = txThank.text_frame
    tf.paragraphs[0].text = "Terima Kasih  •  Mohon Masukan & Saran dari Bapak/Ibu Penguji"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_slide_number(slide15, 15)

    # ─── Save ─────────────────────────────────────────────────────────
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "SIMONAS_Seminar_Hasil_TA.pptx")
    prs.save(output_path)
    print(f"[OK] Presentasi berhasil dibuat: {output_path}")
    print(f"     Total slide: {len(prs.slides)} slide")
    return output_path


if __name__ == "__main__":
    create_presentation()
